"""PDF to PPTX conversion using PyMuPDF + python-pptx."""
import io

import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from ..utils.filenames import temp_output


async def convert_to_pptx(input_path: str) -> str:
    doc = fitz.open(input_path)
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for page_num in range(len(doc)):
            page = doc[page_num]
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

            # Render at 200 DPI for crisp slides on modern displays.
            pix = page.get_pixmap(dpi=200)

            # JPEG for smaller files (slides don't need transparency).
            img_data = io.BytesIO()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img.save(img_data, format="JPEG", quality=90, optimize=True)
            img_data.seek(0)

            # Calculate dimensions to fit slide
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            img_ratio = pix.width / pix.height
            slide_ratio = slide_w / slide_h

            if img_ratio > slide_ratio:
                width = slide_w
                height = int(slide_w / img_ratio)
            else:
                height = slide_h
                width = int(slide_h * img_ratio)

            left = int((slide_w - width) / 2)
            top = int((slide_h - height) / 2)

            slide.shapes.add_picture(img_data, left, top, width, height)
    finally:
        doc.close()

    output_path = temp_output("converted", "pptx")
    prs.save(str(output_path))
    return str(output_path)
