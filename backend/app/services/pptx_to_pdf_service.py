import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from ..utils.cleanup import get_temp_path, ensure_temp_dir


def pptx_to_pdf(input_path: str) -> str:
    """Convert .pptx file to PDF using python-pptx + reportlab."""
    ensure_temp_dir()
    output_path = get_temp_path(f"pptx_{uuid.uuid4().hex}.pdf")

    prs = Presentation(input_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Convert EMU to points (1 inch = 914400 EMU, 1 inch = 72 points)
    pdf_width = slide_width / 914400 * 72
    pdf_height = slide_height / 914400 * 72

    c = canvas.Canvas(str(output_path), pagesize=(pdf_width, pdf_height))

    for slide_num, slide in enumerate(prs.slides):
        if slide_num > 0:
            c.showPage()

        # Draw background
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, pdf_width, pdf_height, fill=1, stroke=0)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Convert shape position from EMU to points
            left = shape.left / 914400 * 72 if shape.left else 0
            top = shape.top / 914400 * 72 if shape.top else 0
            sw = shape.width / 914400 * 72 if shape.width else pdf_width
            sh = shape.height / 914400 * 72 if shape.height else 50

            # Y is from bottom in reportlab
            y = pdf_height - top

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    y -= 10
                    continue

                # Detect font size from first run
                font_size = 12
                is_bold = False
                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        font_size = run.font.size.pt
                    is_bold = run.font.bold or False

                font_size = min(font_size, 36)  # Cap for safety

                if is_bold:
                    c.setFont("Helvetica-Bold", font_size)
                else:
                    c.setFont("Helvetica", font_size)

                c.setFillColorRGB(0.1, 0.1, 0.1)

                # Word wrap within shape width
                words = text.split()
                line = ""
                for word in words:
                    test = f"{line} {word}".strip()
                    if c.stringWidth(test, c._fontname, font_size) < sw - 10:
                        line = test
                    else:
                        if line:
                            c.drawString(left + 5, y, line)
                            y -= font_size + 4
                        line = word
                if line:
                    c.drawString(left + 5, y, line)
                    y -= font_size + 4

    c.save()
    return str(output_path)
